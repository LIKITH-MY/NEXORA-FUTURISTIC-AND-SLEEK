# PPT generation logic
from pptx import Presentation
from utils.logger import log_interaction

def generate_ppt_from_text(
        slides_text,
        output_filename="nexora_presentation.pptx"
):

    try:

        prs = Presentation()

        blocks = [
            block.strip()
            for block in slides_text.split("\n\n")
            if block.strip()
        ]

        if not blocks:
            slide = prs.slides.add_slide(prs.slide_layouts[0])
            slide.shapes.title.text = "Empty Presentation"

        else:

            for block in blocks:

                lines = [
                    line.strip()
                    for line in block.split("\n")
                    if line.strip()
                ]

                slide = prs.slides.add_slide(
                    prs.slide_layouts[1]
                )

                slide.shapes.title.text = lines[0]

                if len(lines) > 1:

                    text_frame = slide.placeholders[1].text_frame

                    for line in lines[1:]:

                        p = text_frame.add_paragraph()
                        p.text = line
                        p.level = 0

        prs.save(output_filename)

        log_interaction(
            "ppt_generation",
            slides_text[:200],
            output_filename
        )

        return output_filename, None

    except Exception as e:

        return None, str(e)
