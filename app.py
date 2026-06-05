


Pasted text(2).txt
Document
i need to upload this project to git hub so i need like python html css and everything in different files so that i can upload them easily to github

Your project is currently in one large Python file (nexora_enhanced.py). 


For GitHub, you should split it like this:

NEXORA/
│
├── app.py
├── requirements.txt
├── README.md
│
├── assets/
│   └── style.css
│
├── modules/
│   ├── chat.py
│   ├── youtube.py
│   ├── weather.py
│   ├── image_gen.py
│   ├── content_creator.py
│   ├── ppt_generator.py
│   ├── voice_assistant.py
│   ├── pdf_chat.py
│   └── snake_game.py
│
├── utils/
│   ├── logger.py
│   └── helpers.py
│
└── data/
    └── audit_log.json
app.py
Contains:

Streamlit UI

Sidebar navigation

Calls functions from modules

assets/style.css
Move all CSS from:

st.markdown("""
<style>
...
</style>
""")
into:

/* style.css */
.main{
    background: linear-gradient(...);
}
and load using:

with open("assets/style.css") as f:
    st.markdown(f"<style>{f.read()}</style>", unsafe_allow_html=True)
modules/chat.py
Contains:

def chat_with_ai():
modules/youtube.py
Contains:

def search_youtube():
modules/weather.py
Contains:

def get_weather():
modules/image_gen.py
Contains:

def generate_image():
modules/content_creator.py
Contains:

CONTENT_CATEGORIES
def generate_content_from_template()
modules/ppt_generator.py
Contains:

def generate_ppt_from_text()
modules/voice_assistant.py
Contains:

def speak_text()
def capture_voice()
modules/pdf_chat.py
Contains:

def extract_text_from_pdf()
modules/snake_game.py
Contains:

def run_snake_game()
utils/logger.py
Contains:

log_interaction()
requirements.txt
streamlit
openai
requests
pillow
beautifulsoup4
python-pptx
PyPDF2
pyttsx3
SpeechRecognition
pygame
IMPORTANT
Before uploading to GitHub remove these keys from the code:

OPENAI_API_KEY = "..."
WEATHER_API_KEY = "..."
Use:

import os

OPENAI_API_KEY = os.getenv("OPENAI_API_KEY")
WEATHER_API_KEY = os.getenv("WEATHER_API_KEY")
Never upload real API keys to GitHub because anyone can use them.

I can also generate the complete separated files (ready-to-upload GitHub project structure) if you want.



Pasted markdown.md
File

Yes. For GitHub, if your goal is simply to upload and showcase the project, then use one file only:

NEXORA/
│
├── nexora_enhanced.py
├── requirements.txt
├── README.md
├── .gitignore
└── audit_log.json
Rename:

nexora_enhanced.py
to

app.py
Then your app.py is literally the entire code you pasted from:

import os
import json
import datetime
import threading
...
all the way to:

st.markdown("---")
The complete file you pasted is your app.py. 


Before uploading:

Remove API Keys
Replace:

OPENAI_API_KEY = "your_key"
WEATHER_API_KEY = "your_key"
with:

import os

OPENAI_API_KEY = os.getenv("OPENAI_API_KEY")
WEATHER_API_KEY = os.getenv("WEATHER_API_KEY")
requirements.txt
streamlit
openai
requests
pillow
beautifulsoup4
python-pptx
PyPDF2
pyttsx3
SpeechRecognition
pygame
.gitignore
__pycache__/
*.pyc
.env
.vscode/
.idea/
audit_log.json
Run
streamlit run app.py
For a student portfolio project, this is perfectly fine. You do not need to split it into 10 files before uploading. Just upload the single working app.py, README.md, and requirements.txt, and you'll have a clean GitHub repository.


give me all in one

I can't give you the complete all-in-one app.py because the code contained in the uploaded text is incomplete and truncated in several places (for example, the CSS section is cut off, parts of the YouTube parser are missing, and the file ends with "Footer" without the full code). 


If I reconstruct it from what's visible, you'll end up with a broken file that won't run.

What I need is the actual source file:

📎 Upload nexora_enhanced.py itself (the .py file, not pasted text or markdown export).

Once you upload the real Python file, I can:

Extract the entire code correctly.

Generate a complete app.py.

Create a GitHub-ready ZIP with:

app.py

requirements.txt

README.md

.gitignore

with no missing lines.

Right now the pasted document is only a partial copy of the original source, so I can't safely reconstruct the full runnable application. 



then what should i put in app.py

For now, put the entire original Nexora code into app.py.

Your GitHub structure should be:

NEXORA/
│
├── app.py
├── requirements.txt
├── README.md
├── .gitignore
└── audit_log.json
app.py
Copy everything from your original nexora_enhanced.py file into app.py:

import os
import json
import datetime
import threading
import random
import re
from io import BytesIO
from urllib.parse import quote_plus, urlencode

import streamlit as st
import requests
from PIL import Image
from bs4 import BeautifulSoup

...
all the way down to the last line of the file.

Before uploading
Remove the hardcoded keys:

OPENAI_API_KEY = "..."
WEATHER_API_KEY = "..."
Replace with:

import os

OPENAI_API_KEY = os.getenv("OPENAI_API_KEY")
WEATHER_API_KEY = os.getenv("WEATHER_API_KEY")
requirements.txt
streamlit
openai
requests
pillow
beautifulsoup4
python-pptx
PyPDF2
pyttsx3
SpeechRecognition
pygame
Run locally
streamlit run app.py
Since the code you pasted is incomplete in the chat export, I cannot generate a new complete app.py from it. The safest approach is:

Open your original nexora_enhanced.py.

Copy all contents.

Paste into app.py.

Upload to GitHub.

That will work immediately and is perfectly fine for a portfolio project.


nexora_enhanced.py
import os
import json
import datetime
import threading
import random
import re
from io import BytesIO
from urllib.parse import quote_plus, urlencode

Third-party libs
import streamlit as st
import requests
from PIL import Image
from bs4 import BeautifulSoup

Optional libs
try:
from openai import OpenAI
openai_client = None # Will be initialized with API key
except Exception:
OpenAI = None
openai_client = None

try:
from pptx import Presentation
from pptx.util import Pt
except Exception:
Presentation = None

try:
from PyPDF2 import PdfReader
except Exception:
PdfReader = None

try:
import pyttsx3
except Exception:
pyttsx3 = None

try:
import speech_recognition as sr
except Exception:
sr = None

try:
import pygame
except Exception:
pygame = None

-------------------------
CONFIG
-------------------------
st.set_page_config(
page_title="NEXORA AI Assistant",
page_icon="🤖",
layout="wide",
initial_sidebar_state="expanded"
)

Custom CSS for attractive UI
st.markdown("""

""", unsafe_allow_html=True)

-------------------------
API KEYS
-------------------------
OPENAI_API_KEY = ""
WEATHER_API_KEY = ""

Initialize OpenAI client with new API
if OpenAI is not None:
try:
openai_client = OpenAI(api_key=OPENAI_API_KEY)
except Exception as e:
openai_client = None
print(f"Failed to initialize OpenAI client: {e}")

-------------------------
CONTENT CATEGORIES
-------------------------
CONTENT_CATEGORIES = {
"Technology": {
"icon": "💻",
"description": "Tech articles, coding, AI & innovation",
"templates": {
"Tech Blog Post": "Write a comprehensive blog post about [TOPIC] covering current trends, benefits, challenges, and future outlook.",
"Product Review": "Create a detailed tech product review for [PRODUCT] including features, pros, cons, comparison with competitors, and recommendation.",
"Tutorial Guide": "Generate a step-by-step tutorial on [SKILL/TOOL] suitable for beginners with clear explanations and examples.",
"Tech News Summary": "Summarize recent developments in [TECH FIELD] with key insights and implications.",
"Code Documentation": "Create comprehensive documentation for [CODE/PROJECT] including setup, usage, and examples."
}
},
"Business": {
"icon": "💼",
"description": "Marketing, strategy, entrepreneurship",
"templates": {
"Business Plan": "Create a business plan for [BUSINESS IDEA] including executive summary, market analysis, financial projections, and strategy.",
"Marketing Copy": "Write compelling marketing copy for [PRODUCT/SERVICE] that highlights benefits and drives conversions.",
"Email Campaign": "Generate an email marketing campaign for [CAMPAIGN GOAL] with subject lines, body content, and call-to-action.",
"Case Study": "Develop a business case study about [COMPANY/PROJECT] showcasing challenges, solutions, and results.",
"SWOT Analysis": "Perform a comprehensive SWOT analysis for [BUSINESS/PROJECT] identifying strengths, weaknesses, opportunities, and threats."
}
},
"Health & Wellness": {
"icon": "🏥",
"description": "Fitness, nutrition, mental health",
"templates": {
"Workout Plan": "Create a [DURATION] workout plan for [FITNESS GOAL] including exercises, sets, reps, and progression.",
"Meal Plan": "Generate a [DURATION] meal plan for [DIETARY GOAL] with recipes, portions, and nutritional information.",
"Wellness Article": "Write an informative article about [HEALTH TOPIC] with evidence-based information and practical tips.",
"Meditation Guide": "Create a guided meditation script for [PURPOSE] with breathing exercises and visualization.",
"Health Tips": "Provide actionable health tips for [SPECIFIC CONDITION/GOAL] backed by research."
}
},
"Education": {
"icon": "📚",
"description": "Learning materials, courses, tutoring",
"templates": {
"Lesson Plan": "Develop a comprehensive lesson plan for teaching [SUBJECT/TOPIC] including objectives, activities, and assessments.",
"Study Guide": "Create a study guide for [EXAM/SUBJECT] with key concepts, practice questions, and study strategies.",
"Course Outline": "Generate a complete course outline for [COURSE TOPIC] with modules, learning objectives, and materials.",
"Educational Article": "Write an educational article explaining [CONCEPT] in simple terms with examples and analogies.",
"Quiz Questions": "Create [NUMBER] quiz questions about [TOPIC] with multiple choice answers and explanations."
}
},
"Entertainment": {
"icon": "🎬",
"description": "Stories, scripts, creative content",
"templates": {
"Short Story": "Write a creative short story about [THEME/PLOT] with compelling characters and an engaging narrative.",
"Movie Script Scene": "Create a movie script scene for [SCENARIO] with dialogue, action, and emotional depth.",
"Blog Entertainment": "Write an entertaining blog post about [TOPIC] with humor, anecdotes, and engaging storytelling.",
"Social Media Content": "Generate engaging social media posts about [TOPIC] optimized for [PLATFORM].",
"Podcast Script": "Create a podcast episode script about [TOPIC] with introduction, main content, and conclusion."
}
},
"Travel": {
"icon": "✈️",
"description": "Guides, itineraries, travel tips",
"templates": {
"Travel Guide": "Create a comprehensive travel guide for [DESTINATION] including attractions, accommodation, food, and tips.",
"Itinerary": "Generate a [DURATION] day itinerary for [DESTINATION] with activities, timing, and recommendations.",
"Travel Blog": "Write a travel blog post about [EXPERIENCE/DESTINATION] with personal insights and practical information.",
"Packing List": "Create a detailed packing list for [TRIP TYPE/DESTINATION] organized by category.",
"Budget Planning": "Develop a travel budget plan for [DESTINATION] including costs for accommodation, food, activities, and transportation."
}
},
"Food & Recipes": {
"icon": "🍳",
"description": "Recipes, cooking tips, culinary content",
"templates": {
"Recipe": "Create a detailed recipe for [DISH] including ingredients, instructions, prep time, and serving suggestions.",
"Cooking Guide": "Write a comprehensive cooking guide for [TECHNIQUE/CUISINE] with tips and best practices.",
"Food Blog Post": "Create a food blog post about [TOPIC] with personal stories, recipes, and beautiful descriptions.",
"Meal Prep Guide": "Generate a meal prep guide for [GOAL/DIETARY NEED] with shopping list and preparation steps.",
"Restaurant Review": "Write a detailed restaurant review for [RESTAURANT/CUISINE] covering ambiance, food quality, and value."
}
},
"Finance": {
"icon": "💰",
"description": "Investment, budgeting, financial advice",
"templates": {
"Budget Plan": "Create a personal budget plan for [INCOME LEVEL/GOALS] with expense categories and saving strategies.",
"Investment Guide": "Write an investment guide for [INVESTMENT TYPE] covering basics, strategies, and risk management.",
"Financial Article": "Create an informative article about [FINANCIAL TOPIC] with practical advice and examples.",
"Savings Strategy": "Develop a savings strategy for [GOAL] with timeline, milestones, and action steps.",
"Financial Report": "Generate a financial analysis report for [COMPANY/PROJECT] with key metrics and insights."
}
}
}

-------------------------
Helper Functions
-------------------------
AUDIT_FILE = "audit_log.json"
if not os.path.exists(AUDIT_FILE):
with open(AUDIT_FILE, "w", encoding="utf-8") as f:
json.dump([], f, indent=2)

def log_interaction(action, query, response):
try:
with open(AUDIT_FILE, "r", encoding="utf-8") as f:
data = json.load(f)
except Exception:
data = []
entry = {
"timestamp": datetime.datetime.now().strftime("%Y-%m-%d %H:%M:%S"),
"action": action,
"query": query,
"response": response[:500]
}
data.append(entry)
with open(AUDIT_FILE, "w", encoding="utf-8") as f:
json.dump(data, f, indent=2, ensure_ascii=False)

def search_youtube(query, max_results=10):
"""Search YouTube without API key using web scraping"""
try:
search_query = quote_plus(query)
url = f"https://www.youtube.com/results?search_query={search_query}"

    headers = {
        'User-Agent': 'Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36 (KHTML, like Gecko) Chrome/91.0.4472.124 Safari/537.36'
    }
    
    response = requests.get(url, headers=headers, timeout=10)
    
    if response.status_code != 200:
        return []
    
    # Parse the HTML
    soup = BeautifulSoup(response.text, 'html.parser')
    
    # Find script tags containing video data
    scripts = soup.find_all('script')
    videos = []
    
    for script in scripts:
        if 'var ytInitialData = ' in str(script):
            # Extract JSON data
            script_text = str(script)
            start = script_text.find('var ytInitialData = ') + len('var ytInitialData = ')
            end = script_text.find(';</script>', start)
            
            if end == -1:
                continue
                
            json_str = script_text[start:end]
            
            try:
                data = json.loads(json_str)
                
                # Navigate through the JSON structure
                contents = data.get('contents', {}).get('twoColumnSearchResultsRenderer', {}).get('primaryContents', {}).get('sectionListRenderer', {}).get('contents', [])
                
                for content in contents:
                    items = content.get('itemSectionRenderer', {}).get('contents', [])
                    
                    for item in items:
                        video_renderer = item.get('videoRenderer', {})
                        
                        if video_renderer:
                            video_id = video_renderer.get('videoId', '')
                            title = video_renderer.get('title', {}).get('runs', [{}])[0].get('text', 'No Title')
                            
                            # Get channel name
                            channel_name = video_renderer.get('ownerText', {}).get('runs', [{}])[0].get('text', 'Unknown Channel')
                            
                            # Get view count and publish time
                            view_count = video_renderer.get('viewCountText', {}).get('simpleText', 'No views')
                            publish_time = video_renderer.get('publishedTimeText', {}).get('simpleText', 'Unknown')
                            
                            # Get thumbnail
                            thumbnails = video_renderer.get('thumbnail', {}).get('thumbnails', [])
                            thumbnail_url = thumbnails[-1].get('url', '') if thumbnails else ''
                            
                            # Get duration
                            duration = video_renderer.get('lengthText', {}).get('simpleText', 'Unknown')
                            
                            if video_id:
                                videos.append({
                                    'video_id': video_id,
                                    'title': title,
                                    'channel': channel_name,
                                    'views': view_count,
                                    'published': publish_time,
                                    'duration': duration,
                                    'thumbnail': thumbnail_url,
                                    'url': f'https://www.youtube.com/watch?v={video_id}'
                                })
                                
                                if len(videos) >= max_results:
                                    break
                    
                    if len(videos) >= max_results:
                        break
                
                if len(videos) >= max_results:
                    break
                    
            except json.JSONDecodeError:
                continue
    
    log_interaction("youtube_search", query, f"Found {len(videos)} videos")
    return videos[:max_results]
    
except Exception as e:
    log_interaction("youtube_search_error", query, str(e))
    return []
def speak_text(text):
if pyttsx3 is None:
return
def _run():
try:
engine = pyttsx3.init()
engine.say(text)
engine.runAndWait()
except Exception:
pass
threading.Thread(target=_run, daemon=True).start()

def capture_voice(timeout=6, phrase_time_limit=7):
if sr is None:
return None, "speech_recognition not available"
r = sr.Recognizer()
try:
with sr.Microphone() as source:
r.adjust_for_ambient_noise(source, duration=0.6)
audio = r.listen(source, timeout=timeout, phrase_time_limit=phrase_time_limit)
text = r.recognize_google(audio)
return text, None
except sr.UnknownValueError:
return None, "Could not understand audio"
except sr.RequestError as e:
return None, f"Recognizer request failed: {e}"
except Exception as e:
return None, f"Voice capture failed: {e}"

def extract_text_from_pdf(uploaded_file) -> str:
if PdfReader is None:
return ""
try:
reader = PdfReader(uploaded_file)
pages = []
for p in reader.pages:
t = p.extract_text()
if t:
pages.append(t)
return "\n".join(pages)
except Exception:
return ""

def chat_with_ai(prompt, context_text=""):
"""Updated to use new OpenAI API"""
if openai_client is None:
return "OpenAI client not available. Please check your API key."
try:
messages = [
{"role": "system", "content": "You are NEXORA, a helpful and intelligent AI assistant."}
]
if context_text:
messages.append({"role": "system", "content": f"Context: {context_text[:1500]}"})
messages.append({"role": "user", "content": prompt})

    # NEW API SYNTAX
    response = openai_client.chat.completions.create(
        model="gpt-4",
        messages=messages,
        max_tokens=800,
        temperature=0.7,
    )
    
    # Access response with new structure
    ans = response.choices[0].message.content.strip()
    log_interaction("chat", prompt, ans)
    return ans
except Exception as e:
    err = f"OpenAI request failed: {e}"
    log_interaction("chat_error", prompt, err)
    return err
def generate_content_from_template(template, user_input):
"""Generate content based on template and user input - Updated API"""
if openai_client is None:
return "OpenAI client not available. Please check your API key."

prompt = template.replace("[TOPIC]", user_input).replace("[PRODUCT]", user_input).replace(
    "[SKILL/TOOL]", user_input).replace("[TECH FIELD]", user_input).replace(
    "[CODE/PROJECT]", user_input).replace("[BUSINESS IDEA]", user_input).replace(
    "[PRODUCT/SERVICE]", user_input).replace("[CAMPAIGN GOAL]", user_input).replace(
    "[COMPANY/PROJECT]", user_input).replace("[BUSINESS/PROJECT]", user_input).replace(
    "[DURATION]", user_input).replace("[FITNESS GOAL]", user_input).replace(
    "[DIETARY GOAL]", user_input).replace("[HEALTH TOPIC]", user_input).replace(
    "[PURPOSE]", user_input).replace("[SPECIFIC CONDITION/GOAL]", user_input).replace(
    "[SUBJECT/TOPIC]", user_input).replace("[EXAM/SUBJECT]", user_input).replace(
    "[COURSE TOPIC]", user_input).replace("[CONCEPT]", user_input).replace(
    "[NUMBER]", user_input).replace("[THEME/PLOT]", user_input).replace(
    "[SCENARIO]", user_input).replace("[PLATFORM]", user_input).replace(
    "[DESTINATION]", user_input).replace("[EXPERIENCE/DESTINATION]", user_input).replace(
    "[TRIP TYPE/DESTINATION]", user_input).replace("[DISH]", user_input).replace(
    "[TECHNIQUE/CUISINE]", user_input).replace("[GOAL/DIETARY NEED]", user_input).replace(
    "[RESTAURANT/CUISINE]", user_input).replace("[INCOME LEVEL/GOALS]", user_input).replace(
    "[INVESTMENT TYPE]", user_input).replace("[FINANCIAL TOPIC]", user_input).replace(
    "[GOAL]", user_input)

try:
    messages = [
        {"role": "system", "content": "You are NEXORA, an expert content creator. Generate high-quality, detailed, and engaging content based on the user's request."},
        {"role": "user", "content": prompt}
    ]
    
    # NEW API SYNTAX
    response = openai_client.chat.completions.create(
        model="gpt-4",
        messages=messages,
        max_tokens=1500,
        temperature=0.8,
    )
    
    content = response.choices[0].message.content.strip()
    log_interaction("content_generation", prompt[:200], content[:200])
    return content
except Exception as e:
    err = f"Content generation failed: {e}"
    log_interaction("content_gen_error", prompt[:200], err)
    return err
def generate_image(prompt):
"""Updated to use new OpenAI API for image generation"""
if openai_client is None:
return None, "OpenAI client not available"
try:
# NEW API SYNTAX
response = openai_client.images.generate(
model="dall-e-2",
prompt=prompt,
n=1,
size="512x512"
)

    url = response.data[0].url
    r = requests.get(url, timeout=20)
    img = Image.open(BytesIO(r.content))
    log_interaction("image_gen", prompt, url)
    return img, None
except Exception as e:
    log_interaction("image_gen_error", prompt, str(e))
    return None, str(e)
def get_weather(city):
if not WEATHER_API_KEY or WEATHER_API_KEY == "YOUR_OPENWEATHER_API_KEY_HERE":
return "⚠️ Weather API key not configured. Please add your OpenWeather API key."
try:
url = f"http://api.openweathermap.org/data/2.5/weather?q={city}&appid={WEATHER_API_KEY}&units=metric"
r = requests.get(url, timeout=10)
d = r.json()
if d.get("cod") != 200:
msg = d.get("message", "City not found")
log_interaction("weather_error", city, msg)
return f"❌ Error: {msg}"
weather = d["weather"][0]["description"].title()
temp = d["main"]["temp"]
hum = d["main"]["humidity"]
wind = d["wind"]["speed"]
out = f"🌍 {city.title()}\n\n🌤️ {weather}\n🌡️ Temperature: {temp}°C\n💧 Humidity: {hum}%\n💨 Wind Speed: {wind} m/s"
log_interaction("weather", city, out)
return out
except Exception as e:
log_interaction("weather_error", city, str(e))
return f"❌ Weather fetch failed: {e}"

def generate_ppt_from_text(slides_text, output_filename="nexora_presentation.pptx"):
if Presentation is None:
return None, "python-pptx not installed"
prs = Presentation()
blocks = [b.strip() for b in slides_text.split("\n\n") if b.strip()]
if not blocks:
slide = prs.slides.add_slide(prs.slide_layouts[1] if len(prs.slide_layouts)>1 else prs.slide_layouts[0])
slide.shapes.title.text = "Empty Slide"
else:
for block in blocks:
lines = [ln.strip() for ln in block.split("\n") if ln.strip()]
if not lines: continue
layout = prs.slide_layouts[1] if len(prs.slide_layouts)>1 else prs.slide_layouts[0]
slide = prs.slides.add_slide(layout)
slide.shapes.title.text = lines[0][:200]
if len(lines) > 1:
try:
tf = slide.placeholders[1].text_frame
for ln in lines[1:]:
p = tf.add_paragraph()
p.text = ln
p.level = 0
except Exception:
slide.shapes.title.text += "\n" + "\n".join(lines[1:])
prs.save(output_filename)
log_interaction("ppt_generate", slides_text[:200], f"Saved: {output_filename}")
return output_filename, None

def run_snake_game():
if pygame is None:
st.error("🎮 pygame not installed or not supported in this environment.")
return
def _run():
WIDTH, HEIGHT = 400, 400
BLOCK = 20
SPEED = 10
UP, DOWN, LEFT, RIGHT = 0,1,2,3

    class SnakeGame:
        def __init__(self):
            pygame.init()
            self.display = pygame.display.set_mode((WIDTH, HEIGHT))
            pygame.display.set_caption("Snake - Nexora")
            self.clock = pygame.time.Clock()
            self.font = pygame.font.SysFont("arial", 25)
            self.reset()

        def reset(self):
            self.head = [WIDTH//2, HEIGHT//2]
            self.snake = [self.head[:], [self.head[0]-BLOCK,self.head[1]], [self.head[0]-2*BLOCK,self.head[1]]]
            self.direction = RIGHT
            self.food = self.place_food()
            self.score = 0

        def place_food(self):
            while True:
                x = random.randint(0, (WIDTH-BLOCK)//BLOCK)*BLOCK
                y = random.randint(0, (HEIGHT-BLOCK)//BLOCK)*BLOCK
                if [x,y] not in self.snake:
                    return [x,y]

        def play_step(self):
            for event in pygame.event.get():
                if event.type == pygame.QUIT:
                    pygame.quit()
                    return False
                elif event.type == pygame.KEYDOWN:
                    if event.key == pygame.K_UP and self.direction != DOWN:
                        self.direction = UP
                    elif event.key == pygame.K_DOWN and self.direction != UP:
                        self.direction = DOWN
                    elif event.key == pygame.K_LEFT and self.direction != RIGHT:
                        self.direction = LEFT
                    elif event.key == pygame.K_RIGHT and self.direction != LEFT:
                        self.direction = RIGHT

            self.move()
            self.snake.insert(0, self.head[:])

            if self.is_collision():
                self.reset()

            if self.head == self.food:
                self.score += 1
                self.food = self.place_food()
            else:
                self.snake.pop()

            self.update_ui()
            self.clock.tick(SPEED)
            return True

        def move(self):
            if self.direction == RIGHT:
                self.head[0] += BLOCK
            elif self.direction == LEFT:
                self.head[0] -= BLOCK
            elif self.direction == DOWN:
                self.head[1] += BLOCK
            elif self.direction == UP:
                self.head[1] -= BLOCK

        def is_collision(self):
            if self.head[0] >= WIDTH or self.head[0] < 0 or self.head[1] >= HEIGHT or self.head[1] < 0:
                return True
            if self.head in self.snake[1:]:
                return True
            return False

        def update_ui(self):
            self.display.fill((0,0,0))
            for pt in self.snake:
                pygame.draw.rect(self.display, (0,255,0), pygame.Rect(pt[0], pt[1], BLOCK, BLOCK))
            pygame.draw.rect(self.display, (255,0,0), pygame.Rect(self.food[0], self.food[1], BLOCK, BLOCK))
            text = self.font.render(f"Score: {self.score}", True, (255,255,255))
            self.display.blit(text, [0,0])
            pygame.display.flip()

    game = SnakeGame()
    running = True
    while running:
        running = game.play_step()

threading.Thread(target=_run, daemon=True).start()
-------------------------
MAIN UI
-------------------------
Header
st.markdown("""

Sidebar
st.sidebar.markdown("### 🎯 Select Tool")
tool_options = {
"💬 Chat": "Chat",
"📄 PDF Chat": "PDF Chat",
"🎤 Voice Assistant": "Voice Assistant",
"📊 PPT Generator": "PPT Generator",
"🌤️ Weather": "Weather",
"🎨 Image Gen": "Image Gen",
"📝 Content Creator": "Content Creator",
"🎥 YouTube Search": "YouTube Search",
"📋 Audit Trail": "Audit Trail",
"🎮 Play Snake": "Play Snake"
}
tool = st.sidebar.selectbox("", list(tool_options.keys()), label_visibility="collapsed")
tool = tool_options[tool]

Session state
if "last_response" not in st.session_state:
st.session_state["last_response"] = ""
if "chat_history" not in st.session_state:
st.session_state["chat_history"] = []
if "selected_category" not in st.session_state:
st.session_state["selected_category"] = None
if "selected_template" not in st.session_state:
st.session_state["selected_template"] = None
if "generated_content" not in st.session_state:
st.session_state["generated_content"] = ""
if "youtube_results" not in st.session_state:
st.session_state["youtube_results"] = []

-------------------------
CHAT
-------------------------
if tool == "Chat":
col1, col2 = st.columns([3, 1])

with col1:
    st.markdown("### 💬 Chat with NEXORA AI")
with col2:
    if st.button("🔄 Clear Chat"):
        st.session_state["chat_history"] = []
        st.rerun()

# Display chat history
if st.session_state["chat_history"]:
    for msg in st.session_state["chat_history"]:
        if msg["role"] == "user":
            st.markdown(f'<div class="chat-bubble user-bubble">👤 {msg["content"]}</div>', unsafe_allow_html=True)
        else:
            st.markdown(f'<div class="chat-bubble ai-bubble">🤖 {msg["content"]}</div>', unsafe_allow_html=True)

# Input
prompt = st.text_input("💭 Ask me anything...", key="chat_input")

col1, col2 = st.columns([1, 1])
with col1:
    if st.button("🚀 Send Message", use_container_width=True):
        if openai_client is None:
            st.error("❌ OpenAI API key missing or openai lib not installed.")
        elif prompt.strip() == "":
            st.info("💡 Please enter a question.")
        else:
            with st.spinner("🤔 Thinking..."):
                resp = chat_with_ai(prompt)
                st.session_state["chat_history"].append({"role": "user", "content": prompt})
                st.session_state["chat_history"].append({"role": "ai", "content": resp})
                st.session_state["last_response"] = resp
                speak_text(resp)
                st.rerun()

with col2:
    if st.button("🔊 Read Last Response", use_container_width=True):
        if st.session_state["last_response"]:
            speak_text(st.session_state["last_response"])
            st.success("🎵 Playing audio...")
        else:
            st.info("No response to read yet.")
-------------------------
PDF CHAT
-------------------------
elif tool == "PDF Chat":
st.markdown("### 📄 PDF-Powered Chat")

uploaded = st.file_uploader("📎 Upload PDF for context", type=["pdf"])
pdf_text = ""

if uploaded:
    with st.spinner("📖 Extracting PDF text..."):
        pdf_text = extract_text_from_pdf(uploaded)
        if pdf_text:
            st.success(f"✅ PDF loaded — {len(pdf_text)} characters extracted.")
            with st.expander("👀 Preview PDF Content"):
                st.text_area("", pdf_text[:1000] + "...", height=200, disabled=True)
        else:
            st.warning("⚠️ No text extracted from PDF.")

user_q = st.text_input("❓ Ask about the document:")

if st.button("🔍 Ask AI", use_container_width=True):
    if user_q.strip() == "":
        st.info("💡 Please type a question.")
    else:
        with st.spinner("🤔 Analyzing document..."):
            context_preview = pdf_text[:3000]
            resp = chat_with_ai(user_q, context_text=context_preview)
            st.session_state["last_response"] = resp
            
            st.markdown('<div class="response-container">', unsafe_allow_html=True)
            st.markdown("**🤖 AI Response:**")
            st.write(resp)
            st.markdown('</div>', unsafe_allow_html=True)
            
            speak_text(resp)
-------------------------
VOICE ASSISTANT
-------------------------
elif tool == "Voice Assistant":
st.markdown("### 🎤 Voice Assistant")
st.info("🎙️ Click 'Record' to speak. Works best on local desktop environments.")

col1, col2, col3 = st.columns(3)

with col1:
    if st.button("🎙️ Record Voice", use_container_width=True):
        with st.spinner("🎧 Listening..."):
            text, err = capture_voice()
            if err:
                st.error(f"❌ {err}")
            else:
                st.success(f"✅ Recognized: **{text}**")
                st.session_state["voice_input"] = text

with col2:
    if st.button("🤖 Ask AI", use_container_width=True):
        voice_txt = st.session_state.get("voice_input", "")
        if not voice_txt:
            st.info("💡 Record voice first.")
        else:
            with st.spinner("🤔 Processing..."):
                resp = chat_with_ai(voice_txt)
                st.session_state["last_response"] = resp
                
                st.markdown('<div class="response-container">', unsafe_allow_html=True)
                st.write(resp)
                st.markdown('</div>', unsafe_allow_html=True)
                
                speak_text(resp)

with col3:
    if st.button("🔊 Read Response", use_container_width=True):
        if st.session_state["last_response"]:
            speak_text(st.session_state["last_response"])
            st.success("🎵 Playing...")
        else:
            st.info("No response yet.")
-------------------------
PPT GENERATOR
-------------------------
elif tool == "PPT Generator":
st.markdown("### 📊 PowerPoint Generator")
st.info("💡 Separate slides with blank lines. First line = title, rest = bullets.")

slides_input = st.text_area(
    "Enter slide content:",
    height=300,
    placeholder="Slide 1 Title\n- Point 1\n- Point 2\n\nSlide 2 Title\n- Point 1"
)

if st.button("✨ Generate PPT", use_container_width=True):
    if not slides_input.strip():
        st.error("❌ Please enter slide content.")
    else:
        with st.spinner("🎨 Creating presentation..."):
            fname, err = generate_ppt_from_text(slides_input)
            if err:
                st.error(f"❌ {err}")
            else:
                with open(fname, "rb") as f:
                    st.download_button(
                        "📥 Download PPT",
                        data=f,
                        file_name=fname,
                        mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                        use_container_width=True
                    )
                st.success(f"✅ PPT generated: {fname}")
-------------------------
WEATHER
-------------------------
elif tool == "Weather":
st.markdown("### 🌤️ Weather Forecast")

col1, col2 = st.columns([3, 1])

with col1:
    city = st.text_input("🌍 Enter city name:", placeholder="e.g., London")

with col2:
    st.write("")
    st.write("")
    if st.button("🔍 Get Weather", use_container_width=True):
        if not city.strip():
            st.info("💡 Please enter a city name.")
        else:
            with st.spinner("🌐 Fetching weather data..."):
                result = get_weather(city.strip())
                
                st.markdown('<div class="response-container">', unsafe_allow_html=True)
                st.markdown(result)
                st.markdown('</div>', unsafe_allow_html=True)
-------------------------
IMAGE GEN
-------------------------
elif tool == "Image Gen":
st.markdown("### 🎨 AI Image Generator")

img_prompt = st.text_area(
    "🖼️ Describe your image:",
    height=100,
    placeholder="A futuristic city at sunset with flying cars..."
)

if st.button("✨ Generate Image", use_container_width=True):
    if openai_client is None:
        st.error("❌ OpenAI API key missing or openai lib not installed.")
    elif not img_prompt.strip():
        st.info("💡 Please enter an image prompt.")
    else:
        with st.spinner("🎨 Creating your image..."):
            img, err = generate_image(img_prompt)
            if img:
                st.image(img, use_column_width=True, caption="Generated Image")
                st.success("✅ Image generated successfully!")
            else:
                st.error(f"❌ Generation failed: {err}")
-------------------------
CONTENT CREATOR
-------------------------
elif tool == "Content Creator":
st.markdown("### 📝 Category-Based Content Creator")
st.info("💡 Select a category and template to generate professional content instantly!")

# Category Selection
if st.session_state["selected_category"] is None:
    st.markdown("#### 🎯 Choose a Content Category")
    
    cols = st.columns(3)
    categories = list(CONTENT_CATEGORIES.keys())
    
    for idx, category in enumerate(categories):
        col_idx = idx % 3
        with cols[col_idx]:
            cat_info = CONTENT_CATEGORIES[category]
            if st.button(
                f"{cat_info['icon']} {category}",
                key=f"cat_{category}",
                use_container_width=True
            ):
                st.session_state["selected_category"] = category
                st.session_state["selected_template"] = None
                st.rerun()
            
            st.markdown(f"""
            <div style="background: white; padding: 1rem; border-radius: 10px; margin-bottom: 1rem; text-align: center; min-height: 100px;">
                <div style="font-size: 2.5rem;">{cat_info['icon']}</div>
                <div style="font-weight: 600; margin: 0.5rem 0;">{category}</div>
                <div style="font-size: 0.85rem; color: #666;">{cat_info['description']}</div>
            </div>
            """, unsafe_allow_html=True)

# Template Selection
elif st.session_state["selected_template"] is None:
    category = st.session_state["selected_category"]
    cat_info = CONTENT_CATEGORIES[category]
    
    col1, col2 = st.columns([3, 1])
    with col1:
        st.markdown(f"#### {cat_info['icon']} {category} Templates")
    with col2:
        if st.button("⬅️ Back", use_container_width=True):
            st.session_state["selected_category"] = None
            st.rerun()
    
    st.markdown(f"*{cat_info['description']}*")
    st.markdown("---")
    
    for template_name, template_desc in cat_info['templates'].items():
        st.markdown(f"""
        <div class="template-card">
            <div class="template-title">📄 {template_name}</div>
            <div class="template-desc">{template_desc[:100]}...</div>
        </div>
        """, unsafe_allow_html=True)
        
        if st.button(f"Use '{template_name}' Template", key=f"tmpl_{template_name}", use_container_width=True):
            st.session_state["selected_template"] = template_name
            st.session_state["generated_content"] = ""
            st.rerun()

# Content Generation
else:
    category = st.session_state["selected_category"]
    template_name = st.session_state["selected_template"]
    cat_info = CONTENT_CATEGORIES[category]
    template_desc = cat_info['templates'][template_name]
    
    col1, col2 = st.columns([3, 1])
    with col1:
        st.markdown(f"#### {cat_info['icon']} {template_name}")
    with col2:
        if st.button("⬅️ Back", use_container_width=True):
            st.session_state["selected_template"] = None
            st.session_state["generated_content"] = ""
            st.rerun()
    
    st.markdown(f"**Template:** {template_desc}")
    st.markdown("---")
    
    # User input for template variables
    st.markdown("##### 📝 Provide Details")
    user_input = st.text_area(
        "Enter your specific topic/details:",
        height=100,
        placeholder="e.g., Artificial Intelligence in Healthcare, My New Fitness App, etc.",
        key="content_input"
    )
    
    col1, col2, col3 = st.columns([1, 1, 1])
    
    with col1:
        if st.button("✨ Generate Content", use_container_width=True):
            if not user_input.strip():
                st.error("❌ Please provide details for the content.")
            elif openai_client is None:
                st.error("❌ OpenAI API key missing or openai lib not installed.")
            else:
                with st.spinner("✍️ Creating your content..."):
                    content = generate_content_from_template(template_desc, user_input)
                    st.session_state["generated_content"] = content
                    st.rerun()
    
    with col2:
        if st.button("🔄 Clear", use_container_width=True):
            st.session_state["generated_content"] = ""
            st.rerun()
    
    with col3:
        if st.button("🏠 New Content", use_container_width=True):
            st.session_state["selected_category"] = None
            st.session_state["selected_template"] = None
            st.session_state["generated_content"] = ""
            st.rerun()
    
    # Display Generated Content
    if st.session_state["generated_content"]:
        st.markdown("---")
        st.markdown("### 📄 Generated Content")
        
        st.markdown('<div class="response-container">', unsafe_allow_html=True)
        st.markdown(st.session_state["generated_content"])
        st.markdown('</div>', unsafe_allow_html=True)
        
        # Download options
        col1, col2, col3 = st.columns(3)
        
        with col1:
            st.download_button(
                "📥 Download as TXT",
                data=st.session_state["generated_content"],
                file_name=f"{template_name.lower().replace(' ', '_')}.txt",
                mime="text/plain",
                use_container_width=True
            )
        
        with col2:
            if st.button("🔊 Read Aloud", use_container_width=True):
                speak_text(st.session_state["generated_content"][:500])
                st.success("🎵 Playing audio...")
        
        with col3:
            if st.button("📋 Copy to Chat", use_container_width=True):
                st.session_state["chat_history"].append({
                    "role": "ai",
                    "content": st.session_state["generated_content"]
                })
                st.success("✅ Added to chat history!")
-------------------------
YOUTUBE SEARCH
-------------------------
elif tool == "YouTube Search":
st.markdown("### 🎥 YouTube Video Search")
st.info("🔍 Search for YouTube videos without needing an API key! Get direct links and video details.")

col1, col2 = st.columns([3, 1])

with col1:
    search_query = st.text_input(
        "🔎 Enter your search query:",
        placeholder="e.g., Python tutorial, cooking recipes, music videos..."
    )

with col2:
    st.write("")
    st.write("")
    num_results = st.selectbox("Results", [5, 10, 15, 20], index=1)

col1, col2 = st.columns([1, 1])

with col1:
    if st.button("🚀 Search YouTube", use_container_width=True):
        if not search_query.strip():
            st.info("💡 Please enter a search query.")
        else:
            with st.spinner("🔍 Searching YouTube..."):
                results = search_youtube(search_query, max_results=num_results)
                st.session_state["youtube_results"] = results
                
                if results:
                    st.success(f"✅ Found {len(results)} videos!")
                    st.rerun()
                else:
                    st.warning("⚠️ No videos found. Try a different search query.")

with col2:
    if st.button("🔄 Clear Results", use_container_width=True):
        st.session_state["youtube_results"] = []
        st.rerun()

# Display Results
if st.session_state["youtube_results"]:
    st.markdown("---")
    st.markdown(f"### 📺 Search Results ({len(st.session_state['youtube_results'])} videos)")
    
    for idx, video in enumerate(st.session_state["youtube_results"], 1):
        st.markdown(f"""
        <div class="video-card">
            <div class="video-title">{idx}. {video['title']}</div>
            <div class="video-channel">📺 {video['channel']}</div>
            <div class="video-stats">
                👁️ {video['views']} • ⏱️ {video['duration']} • 📅 {video['published']}
            </div>
            <a href="{video['url']}" target="_blank" class="video-link">▶️ Watch on YouTube</a>
        </div>
        """, unsafe_allow_html=True)
        
        # Show thumbnail if available
        if video.get('thumbnail'):
            with st.expander(f"🖼️ Show Thumbnail - {video['title'][:50]}..."):
                try:
                    st.image(video['thumbnail'], use_column_width=True)
                except:
                    st.info("Thumbnail not available")
    
    # Export results
    st.markdown("---")
    st.markdown("### 📥 Export Results")
    
    # Prepare text export
    export_text = f"YouTube Search Results for: {search_query}\n"
    export_text += f"Total Videos: {len(st.session_state['youtube_results'])}\n"
    export_text += "="*50 + "\n\n"
    
    for idx, video in enumerate(st.session_state["youtube_results"], 1):
        export_text += f"{idx}. {video['title']}\n"
        export_text += f"   Channel: {video['channel']}\n"
        export_text += f"   Views: {video['views']} | Duration: {video['duration']} | Published: {video['published']}\n"
        export_text += f"   URL: {video['url']}\n\n"
    
    col1, col2 = st.columns(2)
    
    with col1:
        st.download_button(
            "📥 Download as TXT",
            data=export_text,
            file_name=f"youtube_search_{search_query[:30]}.txt",
            mime="text/plain",
            use_container_width=True
        )
    
    with col2:
        # Prepare JSON export
        json_data = json.dumps(st.session_state["youtube_results"], indent=2)
        st.download_button(
            "📥 Download as JSON",
            data=json_data,
            file_name=f"youtube_search_{search_query[:30]}.json",
            mime="application/json",
            use_container_width=True
        )
-------------------------
AUDIT TRAIL
-------------------------
elif tool == "Audit Trail":
st.markdown("### 📋 Audit Trail")

try:
    with open(AUDIT_FILE, "r", encoding="utf-8") as f:
        logs = json.load(f)
except Exception:
    logs = []

col1, col2 = st.columns([3, 1])
with col1:
    st.markdown(f"**Total Entries:** {len(logs)}")
with col2:
    if st.button("📥 Download Audit", use_container_width=True):
        with open(AUDIT_FILE, "rb") as f:
            st.download_button(
                "Download JSON",
                data=f,
                file_name=AUDIT_FILE,
                mime="application/json",
                use_container_width=True
            )

if logs:
    for entry in logs[-30:][::-1]:
        st.markdown(f"""
        <div class="audit-entry">
            <div class="audit-timestamp">🕐 {entry['timestamp']} — {entry['action'].upper()}</div>
            <div style="margin-top: 0.5rem;">
                <strong>❓ Query:</strong> {entry.get('query', 'N/A')[:150]}
            </div>
            <div style="margin-top: 0.5rem;">
                <strong>💬 Response:</strong> {entry.get('response', 'N/A')[:150]}...
            </div>
        </div>
        """, unsafe_allow_html=True)
else:
    st.info("📭 No audit entries yet.")
-------------------------
PLAY SNAKE
-------------------------
elif tool == "Play Snake":
st.markdown("### 🎮 Snake Game")
st.warning("⚠️ Desktop only - launches a pygame window (not supported on hosted Streamlit).")

col1, col2, col3 = st.columns([1, 1, 1])
with col2:
    if st.button("🚀 Launch Game", use_container_width=True):
        if pygame is None:
            st.error("❌ pygame not installed in this environment.")
        else:
            st.success("🎮 Launching Snake game...")
            run_snake_game()
            st.info("🕹️ Game window opened! Use arrow keys to play.")
Footer
st.markdown("---")
st.markdown("""


Close
